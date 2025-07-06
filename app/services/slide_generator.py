"""
Slide Generator Service
Handles content generation and PPTX file creation
"""
import os
import json
from typing import List, Optional, Dict, Any
from datetime import datetime
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.models.presentation import Presentation, Slide, SlideType
from app.config.themes import Theme, ThemeConfig
from app.config.aspect_ratios import AspectRatio, AspectRatioConfig
from app.interfaces.cache import CacheInterface
from app.interfaces.llm import LLMInterface

class SlideGenerator:
    """Service for generating slides and creating PPTX files"""
    
    def __init__(self, cache_service: CacheInterface, llm_service: LLMInterface):
        self.output_dir = "output"
        self.cache = cache_service
        self.llm = llm_service
        os.makedirs(self.output_dir, exist_ok=True)
    
    async def generate_slides(
        self, 
        topic: str, 
        num_slides: int, 
        custom_content: Optional[str] = None,
        theme: Theme = Theme.MODERN,
        font: str = "Arial",
        colors: Optional[Dict[str, str]] = None
    ) -> List[Slide]:
        """
        Generate slides for a given topic with caching
        """
        # Check cache first
        cache_key_params = {
            'topic': topic,
            'num_slides': num_slides,
            'custom_content': custom_content,
            'theme': theme.value if theme else None,
            'font': font,
            'colors': colors
        }
        
        cached_result = self.cache.get_slide_generation(**cache_key_params)
        if cached_result and "slides" in cached_result:
            # Convert cached data back to Slide objects
            slides = [Slide(**slide_data) for slide_data in cached_result["slides"]]
            return slides
        
        # Generate slides if not in cache
        slides = []
        
        # Generate title slide using OpenAI
        try:
            title, subtitle = await self.llm.generate_title_slide_content(topic, custom_content)
            title_slide = Slide(
                slide_type=SlideType.TITLE,
                title=title,
                content=[subtitle],
                citations=[]
            )
        except Exception as e:
            print(f"Failed to generate title slide with OpenAI, using fallback: {str(e)}")
            # Fallback to simple title
            title_slide = Slide(
                slide_type=SlideType.TITLE,
                title=f"{topic}",
                content=[f"Generated on {datetime.now().strftime('%B %d, %Y')}"],
                citations=[]
            )
        slides.append(title_slide)
        
        # Generate content slides
        remaining_slides = num_slides - 1
        if remaining_slides > 0:
            content_slides = await self._generate_content_slides(
                topic, remaining_slides, custom_content
            )
            slides.extend(content_slides)
        
        # Cache the result
        slides_data = [slide.model_dump() for slide in slides]
        self.cache.set_slide_generation(result={"slides": slides_data}, **cache_key_params)
        
        return slides
    
    async def _generate_content_slides(
        self, 
        topic: str, 
        num_slides: int, 
        custom_content: Optional[str] = None
    ) -> List[Slide]:
        """
        Generate content slides using LLM service
        """
        # Use the LLM service to generate slide content
        slides = await self.llm.generate_slides_content(
            topic=topic,
            num_slides=num_slides,
            custom_content=custom_content
        )
        
        return slides
    
    async def create_pptx(self, presentation: Presentation) -> str:
        """
        Create a PPTX file from a presentation
        """
        # Create a new presentation
        pptx = PPTXPresentation()
        
        # Apply theme and styling with aspect ratio
        self._apply_theme(
            pptx, 
            presentation.theme, 
            presentation.aspect_ratio,
            presentation.custom_width,
            presentation.custom_height
        )
        
        # Create slides
        for slide_data in presentation.slides:
            if slide_data.slide_type == SlideType.TITLE:
                self._create_title_slide(pptx, slide_data, presentation)
            elif slide_data.slide_type == SlideType.BULLET_POINTS:
                self._create_bullet_slide(pptx, slide_data, presentation)
            elif slide_data.slide_type == SlideType.TWO_COLUMN:
                self._create_two_column_slide(pptx, slide_data, presentation)
            elif slide_data.slide_type == SlideType.CONTENT_WITH_IMAGE:
                self._create_content_with_image_slide(pptx, slide_data, presentation)
        
        # Save the presentation
        filename = f"presentation_{presentation.id}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        pptx.save(filepath)
        
        return filepath
    
    def _apply_theme(self, pptx: PPTXPresentation, theme: Theme, aspect_ratio: AspectRatio = AspectRatio.WIDESCREEN_16_9, custom_width: Optional[float] = None, custom_height: Optional[float] = None):
        """Apply theme to presentation with proper styling and aspect ratio"""
        # Set slide size based on aspect ratio
        if aspect_ratio == AspectRatio.CUSTOM and custom_width and custom_height:
            # Use custom dimensions
            width_inches, height_inches = Inches(custom_width), Inches(custom_height)
        else:
            # Use predefined aspect ratio dimensions
            width_inches, height_inches = AspectRatioConfig.get_inches_dimensions(aspect_ratio)
        
        pptx.slide_width = width_inches
        pptx.slide_height = height_inches
        
        # Store dimensions for use in slide creation
        self._slide_width = width_inches
        self._slide_height = height_inches
        
        # Get latest theme configuration from centralized config
        theme_config = ThemeConfig.get_theme_config(theme)
        self._theme_colors = theme_config.get("colors", {})
        self._theme_font = theme_config.get("font", "Arial")
        
        # Store theme for reference
        self._current_theme = theme
    
    def _apply_font_and_colors(self, shape, presentation: Presentation, is_title: bool = False):
        """Apply font and colors to a shape with proper priority order and alignment"""
        if not shape.text_frame:
            return
            
        # Get latest theme configuration with priority order
        current_theme = getattr(self, '_current_theme', Theme.MODERN)
        theme_colors = ThemeConfig.get_theme_colors(current_theme)
        theme_font = ThemeConfig.get_theme_font(current_theme)
        custom_colors = presentation.colors or {}
        
        # Configure text frame for proper wrapping and alignment
        shape.text_frame.word_wrap = True
        shape.text_frame.auto_size = True
        
        for paragraph in shape.text_frame.paragraphs:
            # Apply font with priority: custom font > theme font > default
            font_to_use = presentation.font or theme_font or 'Arial'
            paragraph.font.name = font_to_use
            
            # Dynamic font sizing based on content length and slide type
            if is_title:
                # For titles, use larger font but scale down if text is long
                base_size = 24
                if len(paragraph.text) > 50:
                    base_size = 20
                elif len(paragraph.text) > 100:
                    base_size = 16
                elif len(paragraph.text) > 150:
                    base_size = 14
                paragraph.font.size = Pt(base_size)
                # Center align titles/headings
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                # For content, use smaller font with better scaling
                base_size = 14
                if len(paragraph.text) > 100:
                    base_size = 12
                elif len(paragraph.text) > 200:
                    base_size = 10
                elif len(paragraph.text) > 300:
                    base_size = 8
                paragraph.font.size = Pt(base_size)
                # Left align content
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Apply colors with priority order: custom colors > theme colors > default
            if is_title:
                # Title color priority: custom primary > theme primary > default
                color = (custom_colors.get('primary') or 
                        theme_colors.get('primary') or 
                        '#2E86AB')
            else:
                # Content color priority: custom text > theme text > default
                color = (custom_colors.get('text') or 
                        theme_colors.get('text') or 
                        '#2C3E50')
            
            if color.startswith('#'):
                try:
                    r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
                    paragraph.font.color.rgb = RGBColor(r, g, b)
                except (ValueError, IndexError):
                    # Fallback to default color if parsing fails
                    paragraph.font.color.rgb = RGBColor(44, 62, 80)  # Default dark gray
    
    def _apply_background(self, slide, presentation: Presentation):
        """Apply background color to slide with latest theme configuration"""
        # Get latest theme configuration with priority order
        current_theme = getattr(self, '_current_theme', Theme.MODERN)
        theme_colors = ThemeConfig.get_theme_colors(current_theme)
        custom_colors = presentation.colors or {}
        
        # Get background color with priority order: custom colors > theme colors > default
        background_color = (custom_colors.get('background') or 
                           theme_colors.get('background') or 
                           '#FFFFFF')  # Default white
        
        if background_color.startswith('#'):
            try:
                r, g, b = int(background_color[1:3], 16), int(background_color[3:5], 16), int(background_color[5:7], 16)
                # Set the background fill color
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
            except (ValueError, IndexError):
                # Fallback to default white background if parsing fails
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    def _add_citations_box(self, slide, citations, presentation: Presentation):
        """Add a citations text box at the bottom of the slide if citations exist"""
        if not citations:
            return
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        # Combine all citations into a single string
        citations_text = "; ".join(citations)
        
        # Calculate dynamic positioning based on slide dimensions
        slide_width = getattr(self, '_slide_width', Inches(10))
        slide_height = getattr(self, '_slide_height', Inches(7.5))
        
        # Convert to float for calculations
        width_inches = float(slide_width.inches)
        height_inches = float(slide_height.inches)
        
        # Position citations box at bottom with margins
        left = Inches(0.5)
        width = Inches(width_inches - 1.0)  # Full width minus margins
        height = Inches(0.5)
        top = Inches(height_inches - 0.8)  # Near bottom with margin
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = citations_text
        # Style: small font, gray color
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        # Optionally, use the presentation's font
        if presentation.font:
            p.font.name = presentation.font

    def _create_title_slide(self, pptx: PPTXPresentation, slide_data: Slide, presentation: Presentation):
        """Create a title slide with aligned headings and overflow prevention"""
        slide_layout = pptx.slide_layouts[6]  # Blank layout for custom positioning
        slide = pptx.slides.add_slide(slide_layout)
        
        # Apply background first
        self._apply_background(slide, presentation)
        
        # Get slide dimensions for positioning
        slide_width = getattr(self, '_slide_width', Inches(10))
        slide_height = getattr(self, '_slide_height', Inches(7.5))
        
        width_inches = float(slide_width.inches)
        height_inches = float(slide_height.inches)
        
        # Calculate centered positioning for both title and subtitle
        # Use same left position for both to ensure alignment
        left_margin = 1.0  # 1 inch from left
        right_margin = 1.0  # 1 inch from right
        text_width = width_inches - left_margin - right_margin
        
        # Position title in upper center area
        title_top = height_inches * 0.25  # 25% from top
        title_height = 1.5  # 1.5 inches height
        
        # Position subtitle below title with same left alignment
        subtitle_top = title_top + title_height + 0.5  # 0.5 inch gap
        subtitle_height = 1.0  # 1 inch height
        
        # Create title text box
        title_box = slide.shapes.add_textbox(
            Inches(left_margin),
            Inches(title_top),
            Inches(text_width),
            Inches(title_height)
        )
        
        # Create subtitle text box with same left alignment
        subtitle_box = slide.shapes.add_textbox(
            Inches(left_margin),
            Inches(subtitle_top),
            Inches(text_width),
            Inches(subtitle_height)
        )
        
        # Configure title text frame
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_frame.auto_size = True
        
        # Configure subtitle text frame
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.word_wrap = True
        subtitle_frame.auto_size = True
        
        # Set text content
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER
        
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = slide_data.content[0] if slide_data.content else ""
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Ensure subtitle takes full width by setting text frame properties
        subtitle_frame.margin_left = 0
        subtitle_frame.margin_right = 0
        
        # Apply styling with center alignment
        self._apply_font_and_colors(title_box, presentation, is_title=True)
        self._apply_font_and_colors(subtitle_box, presentation, is_title=False)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Add citations if any
        self._add_citations_box(slide, slide_data.citations, presentation)
    
    def _create_bullet_slide(self, pptx: PPTXPresentation, slide_data: Slide, presentation: Presentation):
        """Create a bullet points slide with center-aligned headings and left-aligned content"""
        slide_layout = pptx.slide_layouts[6]  # Blank layout for custom positioning
        slide = pptx.slides.add_slide(slide_layout)
        
        # Apply background first
        self._apply_background(slide, presentation)
        
        # Get slide dimensions for positioning
        slide_width = getattr(self, '_slide_width', Inches(10))
        slide_height = getattr(self, '_slide_height', Inches(7.5))
        
        width_inches = float(slide_width.inches)
        height_inches = float(slide_height.inches)
        
        # Calculate margins and positioning
        left_margin = 0.5  # 0.5 inch from left
        right_margin = 0.5  # 0.5 inch from right
        title_width = width_inches - left_margin - right_margin
        
        # Position title at top with full width
        title_top = 0.5  # 0.5 inch from top
        title_height = 1.0  # 1 inch height
        
        # Position content below title
        content_top = title_top + title_height + 0.3  # 0.3 inch gap
        content_height = height_inches - content_top - 1.0  # Leave space for bottom
        
        # Create title text box with full width
        title_box = slide.shapes.add_textbox(
            Inches(left_margin),
            Inches(title_top),
            Inches(title_width),
            Inches(title_height)
        )
        
        # Create content text box
        content_box = slide.shapes.add_textbox(
            Inches(left_margin),
            Inches(content_top),
            Inches(title_width),
            Inches(content_height)
        )
        
        # Configure title text frame
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_frame.auto_size = True
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        
        # Configure content text frame
        content_frame = content_box.text_frame
        content_frame.clear()
        content_frame.word_wrap = True
        content_frame.auto_size = True
        
        # Set title text
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add bullet points
        for i, point in enumerate(slide_data.content):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = point
            p.level = 0
            # Left align bullet points (content)
            p.alignment = PP_ALIGN.LEFT
        
        # Apply styling
        self._apply_font_and_colors(title_box, presentation, is_title=True)
        self._apply_font_and_colors(content_box, presentation, is_title=False)
        
        # Add citations if any
        self._add_citations_box(slide, slide_data.citations, presentation)
    
    def _create_two_column_slide(self, pptx: PPTXPresentation, slide_data: Slide, presentation: Presentation):
        """Create a two-column slide with center-aligned headings and left-aligned content"""
        slide_layout = pptx.slide_layouts[6]  # Blank layout for custom positioning
        slide = pptx.slides.add_slide(slide_layout)
        
        # Apply background first
        self._apply_background(slide, presentation)
        
        # Get slide dimensions for positioning
        slide_width = getattr(self, '_slide_width', Inches(10))
        slide_height = getattr(self, '_slide_height', Inches(7.5))
        
        width_inches = float(slide_width.inches)
        height_inches = float(slide_height.inches)
        
        # Calculate margins and positioning
        left_margin = 0.5  # 0.5 inch from left
        right_margin = 0.5  # 0.5 inch from right
        title_width = width_inches - left_margin - right_margin
        
        # Position title at top with full width
        title_top = 0.5  # 0.5 inch from top
        title_height = 1.0  # 1 inch height
        
        # Position columns below title
        column_top = title_top + title_height + 0.3  # 0.3 inch gap
        column_height = height_inches - column_top - 1.0  # Leave space for bottom
        
        # Create title text box with full width
        title_box = slide.shapes.add_textbox(
            Inches(left_margin),
            Inches(title_top),
            Inches(title_width),
            Inches(title_height)
        )
        
        # Configure title text frame
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_frame.auto_size = True
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        
        # Set title text
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER
        
        # Calculate column dimensions
        column_width = (title_width - 0.5) / 2  # 0.5" gap between columns
        
        # Create two text boxes for columns with dynamic positioning
        left_box = slide.shapes.add_textbox(
            Inches(left_margin), 
            Inches(column_top), 
            Inches(column_width), 
            Inches(column_height)
        )
        right_box = slide.shapes.add_textbox(
            Inches(left_margin + column_width + 0.5), 
            Inches(column_top), 
            Inches(column_width), 
            Inches(column_height)
        )
        
        # Configure text frames for proper wrapping
        left_frame = left_box.text_frame
        left_frame.clear()
        left_frame.word_wrap = True  # Enable word wrapping
        left_frame.auto_size = True  # Auto-size the text box
        
        right_frame = right_box.text_frame
        right_frame.clear()
        right_frame.word_wrap = True  # Enable word wrapping
        right_frame.auto_size = True  # Auto-size the text box
        
        # Separate content into left and right columns
        left_content = []
        right_content = []
        
        for i, content in enumerate(slide_data.content):
            if content.startswith("Column 1:"):
                # Extract content after "Column 1:"
                clean_content = content.replace("Column 1:", "").strip()
                if clean_content:
                    left_content.append(clean_content)
            elif content.startswith("Column 2:"):
                # Extract content after "Column 2:"
                clean_content = content.replace("Column 2:", "").strip()
                if clean_content:
                    right_content.append(clean_content)
            else:
                # If no column prefix, alternate between left and right
                if len(left_content) <= len(right_content):
                    left_content.append(content)
                else:
                    right_content.append(content)
        
        # Add content to left column
        for i, content in enumerate(left_content):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = content
            # Left align and set proper spacing
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)
        
        # Add content to right column
        for i, content in enumerate(right_content):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = content
            # Left align and set proper spacing
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)
        
        # Apply styling with smaller font size for better fit
        self._apply_font_and_colors(title_box, presentation, is_title=True)
        
        # Apply styling to columns using the new priority-based system
        for paragraph in left_frame.paragraphs:
            if hasattr(paragraph, 'font'):
                if presentation.font:
                    paragraph.font.name = presentation.font
                paragraph.font.size = Pt(12)  # Smaller font for better fit
                
                # Apply colors with priority order
                theme_colors = getattr(self, '_theme_colors', {})
                custom_colors = presentation.colors or {}
                color = (custom_colors.get('text') or 
                        theme_colors.get('text') or 
                        '#2C3E50')
                
                if color.startswith('#'):
                    try:
                        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
                        paragraph.font.color.rgb = RGBColor(r, g, b)
                    except (ValueError, IndexError):
                        paragraph.font.color.rgb = RGBColor(44, 62, 80)
        
        for paragraph in right_frame.paragraphs:
            if hasattr(paragraph, 'font'):
                if presentation.font:
                    paragraph.font.name = presentation.font
                paragraph.font.size = Pt(12)  # Smaller font for better fit
                
                # Apply colors with priority order
                theme_colors = getattr(self, '_theme_colors', {})
                custom_colors = presentation.colors or {}
                color = (custom_colors.get('text') or 
                        theme_colors.get('text') or 
                        '#2C3E50')
                
                if color.startswith('#'):
                    try:
                        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
                        paragraph.font.color.rgb = RGBColor(r, g, b)
                    except (ValueError, IndexError):
                        paragraph.font.color.rgb = RGBColor(44, 62, 80)
        
        # Add citations if any
        self._add_citations_box(slide, slide_data.citations, presentation)
    
    def _create_content_with_image_slide(self, pptx: PPTXPresentation, slide_data: Slide, presentation: Presentation):
        """Create a content slide with center-aligned headings and left-aligned content"""
        slide_layout = pptx.slide_layouts[6]  # Blank layout for custom positioning
        slide = pptx.slides.add_slide(slide_layout)
        
        # Apply background first
        self._apply_background(slide, presentation)
        
        # Get slide dimensions for positioning
        slide_width = getattr(self, '_slide_width', Inches(10))
        slide_height = getattr(self, '_slide_height', Inches(7.5))
        
        width_inches = float(slide_width.inches)
        height_inches = float(slide_height.inches)
        
        # Calculate margins and positioning
        left_margin = 0.5  # 0.5 inch from left
        right_margin = 0.5  # 0.5 inch from right
        title_width = width_inches - left_margin - right_margin
        
        # Position title at top with full width
        title_top = 0.5  # 0.5 inch from top
        title_height = 1.0  # 1 inch height
        
        # Position content below title
        content_top = title_top + title_height + 0.3  # 0.3 inch gap
        content_height = height_inches - content_top - 2.0  # Leave space for image and bottom
        
        # Create title text box with full width
        title_box = slide.shapes.add_textbox(
            Inches(left_margin),
            Inches(title_top),
            Inches(title_width),
            Inches(title_height)
        )
        
        # Create content text box
        content_box = slide.shapes.add_textbox(
            Inches(left_margin),
            Inches(content_top),
            Inches(title_width),
            Inches(content_height)
        )
        
        # Configure title text frame
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_frame.auto_size = True
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        
        # Configure content text frame
        content_frame = content_box.text_frame
        content_frame.clear()
        content_frame.word_wrap = True
        content_frame.auto_size = True
        
        # Set title text
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add content points
        for i, point in enumerate(slide_data.content):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            p.text = point
            # Left align content points
            p.alignment = PP_ALIGN.LEFT
        
        # Add image placeholder at bottom center
        if slide_data.image_suggestion:
            # Create a placeholder for the image at bottom center with proper text wrapping
            # Position: center horizontally, near bottom vertically
            img_placeholder = slide.shapes.add_textbox(Inches(3.5), Inches(5), Inches(3), Inches(1.5))
            img_frame = img_placeholder.text_frame
            img_frame.clear()
            img_frame.word_wrap = True  # Enable word wrapping
            img_frame.auto_size = True  # Auto-size the text box
            
            p = img_frame.paragraphs[0]
            p.text = f"[Image: {slide_data.image_suggestion}]"
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(6)
        
        # Apply styling
        self._apply_font_and_colors(title_box, presentation, is_title=True)
        self._apply_font_and_colors(content_box, presentation, is_title=False)
        if slide_data.image_suggestion:
            self._apply_font_and_colors(img_placeholder, presentation, is_title=False)
        
        # Add citations if any
        self._add_citations_box(slide, slide_data.citations, presentation) 